import streamlit as st
import streamlit.components.v1 as components
import datetime
import html
import json
import os
import tempfile
from pathlib import Path
import hashlib
from typing import List, Dict, Any

# Document processing imports (you'll need to install these)
try:
    import PyPDF2
    from docx import Document
    from PIL import Image
    import pytesseract  # For OCR
    from pptx import Presentation
    import markdown
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    st.warning("⚠️ Some document processing libraries are missing. Install: pip install PyPDF2 python-docx pillow pytesseract python-pptx markdown")

# LLM Integration (using requests for API calls)
import requests

st.set_page_config(page_title="Intituas AI", layout="wide", page_icon="🤖")

# =====================
# CONFIGURATION
# =====================
SUPPORTED_FILE_TYPES = {
    "pdf": "📄 PDF Documents",
    "docx": "📝 Word Documents", 
    "pptx": "📊 PowerPoint Presentations",
    "txt": "📄 Text Files",
    "md": "📝 Markdown Files",
    "png": "🖼️ PNG Images",
    "jpg": "🖼️ JPG Images", 
    "jpeg": "🖼️ JPEG Images"
}

# =====================
# SESSION STATE
# =====================
def initialize_session_state():
    defaults = {
        'current_mode': None,
        'notes': [],
        'current_note': None,
        'show_mode_selector': False,
        'processed_documents': {},  # Cache for processed documents
        'gemini_api_key': '',
        'show_transcript': False,
        'selected_document': None
    }
    
    for key, default_value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default_value

initialize_session_state()

# =====================
# DOCUMENT PROCESSING FUNCTIONS
# =====================
class DocumentProcessor:
    @staticmethod
    def get_file_hash(file_content: bytes) -> str:
        """Generate hash for file content to cache processing results"""
        return hashlib.md5(file_content).hexdigest()
    
    @staticmethod
    def extract_text_from_pdf(file_content: bytes) -> str:
        """Extract text from PDF"""
        if not PDF_AVAILABLE:
            return "PDF processing not available. Please install required libraries."
        
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                tmp_file.write(file_content)
                tmp_file.flush()
                
                text = ""
                with open(tmp_file.name, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                
                os.unlink(tmp_file.name)
                return text
        except Exception as e:
            return f"Error processing PDF: {str(e)}"
    
    @staticmethod
    def extract_text_from_docx(file_content: bytes) -> str:
        """Extract text from Word document"""
        if not PDF_AVAILABLE:
            return "DOCX processing not available. Please install required libraries."
        
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                tmp_file.write(file_content)
                tmp_file.flush()
                
                doc = Document(tmp_file.name)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                
                os.unlink(tmp_file.name)
                return text
        except Exception as e:
            return f"Error processing DOCX: {str(e)}"
    
    @staticmethod
    def extract_text_from_pptx(file_content: bytes) -> str:
        """Extract text from PowerPoint presentation"""
        if not PDF_AVAILABLE:
            return "PPTX processing not available. Please install required libraries."
        
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                tmp_file.write(file_content)
                tmp_file.flush()
                
                prs = Presentation(tmp_file.name)
                text = ""
                for slide_num, slide in enumerate(prs.slides, 1):
                    text += f"Slide {slide_num}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    text += "\n"
                
                os.unlink(tmp_file.name)
                return text
        except Exception as e:
            return f"Error processing PPTX: {str(e)}"
    
    @staticmethod
    def extract_text_from_image(file_content: bytes) -> str:
        """Extract text from image using OCR"""
        if not PDF_AVAILABLE:
            return "Image processing not available. Please install required libraries."
        
        try:
            with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
                tmp_file.write(file_content)
                tmp_file.flush()
                
                image = Image.open(tmp_file.name)
                text = pytesseract.image_to_string(image)
                
                os.unlink(tmp_file.name)
                return text
        except Exception as e:
            return f"Error processing image: {str(e)}. Make sure tesseract is installed."
    
    @staticmethod
    def process_markdown(file_content: bytes) -> str:
        """Process markdown file"""
        try:
            text = file_content.decode('utf-8')
            # Convert markdown to plain text (basic approach)
            import re
            # Remove markdown formatting
            text = re.sub(r'#+ ', '', text)  # Headers
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # Bold
            text = re.sub(r'\*(.*?)\*', r'\1', text)  # Italic
            text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)  # Links
            return text
        except Exception as e:
            return f"Error processing markdown: {str(e)}"
    
    @staticmethod
    def process_document(uploaded_file) -> Dict[str, Any]:
        """Main document processing function"""
        file_content = uploaded_file.read()
        file_hash = DocumentProcessor.get_file_hash(file_content)
        
        # Check cache first
        if file_hash in st.session_state.processed_documents:
            return st.session_state.processed_documents[file_hash]
        
        file_extension = uploaded_file.name.lower().split('.')[-1]
        
        # Process based on file type
        if file_extension == 'pdf':
            text = DocumentProcessor.extract_text_from_pdf(file_content)
        elif file_extension == 'docx':
            text = DocumentProcessor.extract_text_from_docx(file_content)
        elif file_extension == 'pptx':
            text = DocumentProcessor.extract_text_from_pptx(file_content)
        elif file_extension in ['png', 'jpg', 'jpeg']:
            text = DocumentProcessor.extract_text_from_image(file_content)
        elif file_extension == 'md':
            text = DocumentProcessor.process_markdown(file_content)
        elif file_extension == 'txt':
            text = file_content.decode('utf-8')
        else:
            text = "Unsupported file type"
        
        # Create document info
        doc_info = {
            'filename': uploaded_file.name,
            'file_type': file_extension,
            'file_size': len(file_content),
            'text_content': text,
            'word_count': len(text.split()),
            'processed_at': datetime.datetime.now().isoformat(),
            'file_hash': file_hash
        }
        
        # Cache the result
        st.session_state.processed_documents[file_hash] = doc_info
        
        return doc_info

# =====================
# AI/LLM INTEGRATION
# =====================
class AIProcessor:
    @staticmethod
    def get_ai_response(prompt: str, context: str = "") -> str:
        """Get AI response using Gemini API or mock response"""
        try:
            # This is a mock implementation
            # In a real app, you'd integrate with Gemini API using the API key
            
            # Mock responses based on prompt type
            if "summarize" in prompt.lower() or "summary" in prompt.lower():
                return f"📝 **Document Summary:**\n\nBased on the uploaded document(s), here are the key points:\n\n• Main topic: {context[:100]}...\n• Key insights extracted from the content\n• Important conclusions and takeaways\n\n*This is a sample response. Integrate with Gemini API for real functionality.*"
            
            elif "question" in prompt.lower() or "?" in prompt:
                return f"🤔 **Answer:**\n\nBased on the document content, here's what I found:\n\n{context[:200]}...\n\n*This is a sample response. The actual response would analyze your document and provide specific answers.*"
            
            else:
                return f"🤖 **AI Response:**\n\nI understand you're asking about: {prompt}\n\nBased on your document content, I can help you with analysis, summarization, and answering questions.\n\n*Note: This is a demo response. Connect to Gemini API for full functionality.*"
                
        except Exception as e:
            return f"Error getting AI response: {str(e)}"
    
    @staticmethod
    def generate_summary(documents: List[Dict]) -> str:
        """Generate summary of uploaded documents"""
        if not documents:
            return "No documents to summarize."
        
        total_words = sum(doc.get('word_count', 0) for doc in documents)
        file_types = list(set(doc.get('file_type', 'unknown') for doc in documents))
        
        summary = f"📊 **Document Analysis Summary**\n\n"
        summary += f"• **Total Documents:** {len(documents)}\n"
        summary += f"• **Total Words:** {total_words:,}\n"
        summary += f"• **File Types:** {', '.join(file_types)}\n\n"
        
        summary += "**Documents:**\n"
        for i, doc in enumerate(documents, 1):
            summary += f"{i}. **{doc.get('filename', 'Unknown')}** ({doc.get('word_count', 0)} words)\n"
        
        # Add sample content analysis
        if documents:
            first_doc = documents[0]
            content_preview = first_doc.get('text_content', '')[:300]
            summary += f"\n**Content Preview:**\n{content_preview}...\n"
        
        return summary

# =====================
# TOOLS FUNCTIONS
# =====================
class ToolsProcessor:
    @staticmethod
    def create_mindmap(documents: List[Dict]) -> Dict:
        """Create a mind map from document content"""
        if not documents:
            return {"center": "No Documents", "nodes": []}
        
        # Extract key topics (mock implementation)
        all_text = " ".join([doc.get('text_content', '') for doc in documents])
        words = all_text.split()
        
        # Simple keyword extraction (in real implementation, use NLP)
        common_words = ['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'should', 'could', 'may', 'might', 'must', 'can', 'shall']
        filtered_words = [word.lower().strip('.,!?";') for word in words if len(word) > 3 and word.lower() not in common_words]
        
        # Get most frequent words as nodes
        from collections import Counter
        word_counts = Counter(filtered_words)
        top_words = [word for word, count in word_counts.most_common(8)]
        
        return {
            "center": f"Analysis of {len(documents)} Documents",
            "nodes": top_words
        }
    
    @staticmethod
    def cluster_documents(documents: List[Dict]) -> Dict:
        """Cluster documents by content similarity"""
        if len(documents) < 2:
            return {"clusters": [{"name": "Single Document", "documents": documents}]}
        
        # Simple clustering by file type (mock implementation)
        clusters = {}
        for doc in documents:
            file_type = doc.get('file_type', 'unknown')
            if file_type not in clusters:
                clusters[file_type] = []
            clusters[file_type].append(doc)
        
        return {
            "clusters": [
                {"name": f"{file_type.upper()} Documents", "documents": docs}
                for file_type, docs in clusters.items()
            ]
        }
    
    @staticmethod
    def generate_prompts(documents: List[Dict]) -> List[str]:
        """Generate suggested prompts based on document content"""
        base_prompts = [
            "What are the main themes discussed in these documents?",
            "Can you summarize the key findings?",
            "What are the most important conclusions?",
            "Are there any contradictions between the documents?",
            "What questions are left unanswered?",
            "What are the practical implications?",
            "How do these documents relate to each other?",
            "What evidence supports the main arguments?"
        ]
        
        if documents:
            doc_types = set(doc.get('file_type', 'document') for doc in documents)
            if 'pdf' in doc_types:
                base_prompts.append("What are the key sections in the PDF documents?")
            if 'pptx' in doc_types:
                base_prompts.append("What are the main points from the presentations?")
            if len(documents) > 1:
                base_prompts.append("Compare and contrast the different documents.")
        
        return base_prompts[:6]  # Return top 6 prompts

# =====================
# CUSTOM CSS (Enhanced)
# =====================
custom_css = """
<style>
    .note-card { 
        padding: 20px; 
        border-radius: 10px; 
        border: 1px solid #ddd; 
        margin: 10px 0; 
        cursor: pointer; 
        transition: all 0.3s ease;
        background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
    }
    .note-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    }
    .mode-selector { 
        padding: 20px; 
        border-radius: 10px; 
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        margin: 20px 0; 
    }
    .summary-container { 
        background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
        color: white;
        border-radius: 10px; 
        padding: 20px; 
        margin: 10px 0; 
    }
    .doc-info-card {
        background: #f8f9fa;
        border-radius: 8px;
        padding: 15px;
        margin: 10px 0;
        border-left: 4px solid #4CAF50;
        cursor: pointer;
        transition: all 0.3s ease;
    }
    .doc-info-card:hover {
        background: #e9ecef;
        border-left-color: #28a745;
    }
    .doc-info-card.selected {
        background: #e3f2fd;
        border-left-color: #2196f3;
    }
    .ai-response {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
    }
    .transcript-container {
        background: #f8f9fa;
        border-radius: 10px;
        padding: 15px;
        max-height: 400px;
        overflow-y: auto;
        border: 1px solid #dee2e6;
    }
    .tool-card {
        background: linear-gradient(135deg, #84fab0 0%, #8fd3f4 100%);
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
        text-align: center;
        cursor: pointer;
        transition: all 0.3s ease;
    }
    .tool-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    }
    .mindmap-node {
        background: white;
        border: 2px solid #667eea;
        border-radius: 20px;
        padding: 8px 16px;
        margin: 5px;
        display: inline-block;
        font-weight: bold;
    }
    .mindmap-center {
        background: #667eea;
        color: white;
        font-size: 1.1em;
    }
    .cluster-container {
        background: #fff3cd;
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
        border: 1px solid #ffeaa7;
    }
    .prompt-button {
        background: linear-gradient(135deg, #a8edea 0%, #fed6e3 100%);
        border: none;
        border-radius: 20px;
        padding: 10px 15px;
        margin: 5px;
        cursor: pointer;
        transition: all 0.3s ease;
    }
    .prompt-button:hover {
        transform: scale(1.05);
    }
    .stProgress .stProgressBar {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
    }
</style>
"""

# =====================
# HELPER FUNCTIONS (Enhanced)
# =====================
def create_new_note(title, mode):
    new_note = {
        'id': len(st.session_state.notes),
        'title': title,
        'mode': mode,
        'date': datetime.datetime.now().strftime("%b %d, %Y"),
        'documents': [],  # Store processed documents
        'view_mode': 'chat',
        'messages': [],
        'summary_generated': False,
        'active_tools': [],
        'mindmap_data': None,
        'cluster_data': None,
        'suggested_prompts': []
    }
    st.session_state.notes.append(new_note)
    st.session_state.current_note = new_note
    st.session_state.show_mode_selector = False

def build_chat_html(messages):
    items = ""
    for msg in messages:
        css_class = "user" if msg["type"] == "user" else "bot"
        content = html.escape(str(msg["content"])).replace("\\n","<br>").replace("\n", "<br>")
        items += f'<div class="msg {css_class}">{content}</div><div class="clear"></div>'
    
    return f"""
    <!doctype html>
    <html>
      <head>
        <meta charset="utf-8" />
        <style>
          html,body {{ height:100%; margin:0; padding:0; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; }}
          .chat-box {{ height:100%; overflow-y:auto; box-sizing:border-box; padding:15px; border-radius:10px; background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%); }}
          .msg {{ margin:10px 0; padding:12px 16px; border-radius:15px; max-width:75%; clear:both; word-wrap:break-word; line-height: 1.4; }}
          .user {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; float:right; text-align:right; }}
          .bot {{ background: white; float:left; text-align:left; box-shadow: 0 2px 5px rgba(0,0,0,0.1); }}
          .clear {{ clear:both; }}
        </style>
      </head>
      <body>
        <div class="chat-box" id="chat">{items}<div id="end"></div></div>
        <script>
          function scrollToBottom() {{
            var chat = document.getElementById("chat");
            if (chat) {{ chat.scrollTop = chat.scrollHeight; }}
          }}
          setTimeout(scrollToBottom, 100);
        </script>
      </body>
    </html>
    """

def show_research_mode(note):
    """Research mode with side-by-side layout like NotebookLM"""
    st.title(f"🔬 {note['title']}")
    st.caption(f"Research Mode • Created {note['date']}")
    
    # Main layout: Sources on left, Chat on right
    source_col, chat_col = st.columns([1, 1])
    
    with source_col:
        st.subheader("📚 Sources")
        
        # File uploader
        uploaded_files = st.file_uploader(
            "Upload your research documents",
            type=list(SUPPORTED_FILE_TYPES.keys()),
            accept_multiple_files=True,
            key=f"research_docs_{note['id']}"
        )
        
        if uploaded_files:
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            for i, uploaded_file in enumerate(uploaded_files):
                # Check if already processed
                if not any(doc['filename'] == uploaded_file.name for doc in note['documents']):
                    status_text.text(f"Processing {uploaded_file.name}...")
                    progress_bar.progress((i + 1) / len(uploaded_files))
                    
                    # Process the document
                    doc_info = DocumentProcessor.process_document(uploaded_file)
                    note['documents'].append(doc_info)
            
            progress_bar.empty()
            status_text.empty()
        
        # Display processed documents
        if note['documents']:
            for i, doc in enumerate(note['documents']):
                is_selected = st.session_state.get('selected_document') == doc['file_hash']
                
                if st.button(
                    f"📄 {doc['filename']}\n{doc['file_type'].upper()} • {doc['word_count']:,} words",
                    key=f"doc_btn_{doc['file_hash']}",
                    use_container_width=True
                ):
                    st.session_state.selected_document = doc['file_hash']
                    st.session_state.show_transcript = True
                    st.rerun()
        
        # Show transcript if document is selected
        if st.session_state.get('show_transcript') and st.session_state.get('selected_document'):
            selected_doc = None
            for doc in note['documents']:
                if doc['file_hash'] == st.session_state.selected_document:
                    selected_doc = doc
                    break
            
            if selected_doc:
                st.subheader(f"📝 Transcript: {selected_doc['filename']}")
                st.markdown(
                    f"""
                    <div class="transcript-container">
                        <pre style="white-space: pre-wrap; font-family: 'Segoe UI', sans-serif; font-size: 14px;">
                        {selected_doc['text_content'][:2000]}{"..." if len(selected_doc['text_content']) > 2000 else ""}
                        </pre>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
    
    with chat_col:
        st.subheader("💬 AI Chat")
        
        # Chat interface
        components.html(build_chat_html(note["messages"]), height=300, scrolling=True)
        
        # Chat input with Enter key support
        chat_container = st.container()
        with chat_container:
            user_input = st.text_input(
                "Ask questions about your research...",
                key=f"research_input_{note['id']}",
                placeholder="e.g., 'What are the main findings?' or 'Compare the methodologies'"
            )
            
            col1, col2 = st.columns([4, 1])
            with col2:
                send_clicked = st.button("Send", key=f"send_research_{note['id']}")
            
            # Handle message sending
            if send_clicked or (user_input and st.session_state.get(f"research_input_{note['id']}_submit", False)):
                if user_input.strip() and note['documents']:
                    # Add user message
                    note['messages'].append({"type": "user", "content": user_input})
                    
                    # Prepare context from documents
                    context = "\n\n".join([doc['text_content'][:1000] for doc in note['documents']])
                    
                    # Get AI response
                    ai_response = AIProcessor.get_ai_response(user_input, context)
                    note['messages'].append({"type": "bot", "content": ai_response})
                    
                    # Clear input and rerun
                    st.session_state[f"research_input_{note['id']}"] = ""
                    st.rerun()
                    
                elif user_input.strip() and not note['documents']:
                    note['messages'].append({"type": "user", "content": user_input})
                    note['messages'].append({"type": "bot", "content": "Please upload research documents first! 📄"})
                    st.session_state[f"research_input_{note['id']}"] = ""
                    st.rerun()

def show_regular_note_content(note):
    """Regular note content for Study and Normal modes"""
    st.title(f"{'📚' if note['mode'] == 'Study' else '📝'} {note['title']}")
    st.caption(f"{note['mode']} Mode • Created {note['date']}")

    # Document upload section
    with st.expander(f"📚 Documents ({len(note['documents'])})", expanded=True):
        uploaded_files = st.file_uploader(
            "Upload your documents for AI analysis",
            type=list(SUPPORTED_FILE_TYPES.keys()),
            accept_multiple_files=True,
            key=f"docs_{note['id']}"
        )
        
        if uploaded_files:
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            for i, uploaded_file in enumerate(uploaded_files):
                # Check if already processed
                if not any(doc['filename'] == uploaded_file.name for doc in note['documents']):
                    status_text.text(f"Processing {uploaded_file.name}...")
                    progress_bar.progress((i + 1) / len(uploaded_files))
                    
                    # Process the document
                    doc_info = DocumentProcessor.process_document(uploaded_file)
                    note['documents'].append(doc_info)
            
            progress_bar.empty()
            status_text.empty()
        
        # Display processed documents
        if note['documents']:
            st.subheader("📄 Processed Documents")
            for doc in note['documents']:
                with st.container():
                    st.markdown(f"""
                    <div class="doc-info-card">
                        <strong>📄 {doc['filename']}</strong><br>
                        Type: {doc['file_type'].upper()} | Size: {doc['file_size']:,} bytes | Words: {doc['word_count']:,}
                    </div>
                    """, unsafe_allow_html=True)

    # View toggle buttons
    col1, col2 = st.columns([1, 4])
    with col1:
        if st.button("💬 Chat", key=f"chat_btn_{note['id']}", use_container_width=True):
            note['view_mode'] = 'chat'
        if st.button("📝 Summary", key=f"summary_btn_{note['id']}", use_container_width=True):
            note['view_mode'] = 'summary'

    # Main content area
    with col2:
        if note['view_mode'] == 'chat':
            # Chat interface
            components.html(build_chat_html(note["messages"]), height=400, scrolling=True)

            # Chat input with Enter support
            input_key = f"chat_input_{note['id']}"
            
            user_input = st.text_input(
                "Ask questions about your documents...", 
                key=input_key, 
                placeholder="e.g., 'Summarize the main points' or 'What are the key findings?'"
            )
            
            col1, col2 = st.columns([4, 1])
            with col2:
                send_clicked = st.button("Send", key=f"send_{note['id']}")
            
            # Handle sending
            if send_clicked or user_input:  # This will trigger on Enter
                if user_input.strip() and note['documents']:
                    # Add user message
                    note['messages'].append({"type": "user", "content": user_input})
                    
                    # Prepare context from documents
                    context = "\n\n".join([doc['text_content'][:1000] for doc in note['documents']])
                    
                    # Get AI response
                    ai_response = AIProcessor.get_ai_response(user_input, context)
                    note['messages'].append({"type": "bot", "content": ai_response})
                    
                    # Clear input and rerun
                    st.session_state[input_key] = ""
                    st.rerun()
                    
                elif user_input.strip() and not note['documents']:
                    note['messages'].append({"type": "user", "content": user_input})
                    note['messages'].append({"type": "bot", "content": "Please upload documents first so I can analyze them and answer your questions! 📄"})
                    st.session_state[input_key] = ""
                    st.rerun()
                    
        else:  # Summary view
            if note['documents']:
                summary = AIProcessor.generate_summary(note['documents'])
                st.markdown(f"""
                <div class="ai-response">
                {summary}
                </div>
                """, unsafe_allow_html=True)
                
                # Generate AI summary button
                if st.button("🤖 Generate AI Summary", key=f"ai_summary_{note['id']}"):
                    with st.spinner("Generating AI summary..."):
                        context = "\n\n".join([doc['text_content'] for doc in note['documents']])
                        ai_summary = AIProcessor.get_ai_response("Please provide a comprehensive summary", context)
                        st.markdown(f"""
                        <div class="ai-response">
                        {ai_summary}
                        </div>
                        """, unsafe_allow_html=True)
            else:
                st.info("📄 Upload documents to generate summaries and analysis!")

    # Tools section
    show_tools_section(note)

def show_tools_section(note):
    """Display tools section with Mind Map, Clustering, and Suggested Prompts"""
    with st.expander("🛠️ AI Tools", expanded=False):
        if not note['documents']:
            st.info("Upload documents to unlock AI tools!")
            return
            
        tool_cols = st.columns(3)
        
        with tool_cols[0]:
            if st.button("🧠 Mind Map", key=f"mindmap_{note['id']}", use_container_width=True):
                if 'mindmap' not in note['active_tools']:
                    note['active_tools'].append('mindmap')
                    note['mindmap_data'] = ToolsProcessor.create_mindmap(note['documents'])
                else:
                    note['active_tools'].remove('mindmap')
        
        with tool_cols[1]:
            if st.button("🧩 Clustering", key=f"cluster_{note['id']}", use_container_width=True):
                if 'clustering' not in note['active_tools']:
                    note['active_tools'].append('clustering')
                    note['cluster_data'] = ToolsProcessor.cluster_documents(note['documents'])
                else:
                    note['active_tools'].remove('clustering')
        
        with tool_cols[2]:
            if st.button("💡 Suggested Prompts", key=f"prompts_{note['id']}", use_container_width=True):
                if 'prompts' not in note['active_tools']:
                    note['active_tools'].append('prompts')
                    note['suggested_prompts'] = ToolsProcessor.generate_prompts(note['documents'])
                else:
                    note['active_tools'].remove('prompts')

        # Display active tools
        if 'mindmap' in note['active_tools'] and note.get('mindmap_data'):
            st.subheader("🧠 Mind Map")
            mindmap = note['mindmap_data']
            
            # Center node
            st.markdown(f"""
            <div style="text-align: center; margin: 20px 0;">
                <div class="mindmap-node mindmap-center">{mindmap['center']}</div>
            </div>
            """, unsafe_allow_html=True)
            
            # Connected nodes
            if mindmap['nodes']:
                nodes_html = "".join([f'<div class="mindmap-node">{node}</div>' for node in mindmap['nodes']])
                st.markdown(f"""
                <div style="text-align: center;">
                    {nodes_html}
                </div>
                """, unsafe_allow_html=True)
        
        if 'clustering' in note['active_tools'] and note.get('cluster_data'):
            st.subheader("🧩 Document Clusters")
            clusters = note['cluster_data']['clusters']
            
            for cluster in clusters:
                st.markdown(f"""
                <div class="cluster-container">
                    <h4>📁 {cluster['name']}</h4>
                    <ul>
                    {"".join([f"<li>📄 {doc['filename']} ({doc['word_count']} words)</li>" for doc in cluster['documents']])}
                    </ul>
                </div>
                """, unsafe_allow_html=True)
        
        if 'prompts' in note['active_tools'] and note.get('suggested_prompts'):
            st.subheader("💡 Suggested Questions")
            st.markdown("Click on any question to ask it:")
            
            cols = st.columns(2)
            for i, prompt in enumerate(note['suggested_prompts']):
                with cols[i % 2]:
                    if st.button(f"❓ {prompt}", key=f"prompt_{note['id']}_{i}", use_container_width=True):
                        # Add the prompt as a user message
                        note['messages'].append({"type": "user", "content": prompt})
                        
                        # Prepare context and get AI response
                        context = "\n\n".join([doc['text_content'][:1000] for doc in note['documents']])
                        ai_response = AIProcessor.get_ai_response(prompt, context)
                        note['messages'].append({"type": "bot", "content": ai_response})
                        
                        # Switch to chat view and rerun
                        note['view_mode'] = 'chat'
                        st.rerun()

# =====================
# SIDEBAR (Enhanced)
# =====================
with st.sidebar:
    st.markdown("# 🤖 Intituas AI")
    st.markdown("*Powered by Gemini*")
    
    if st.button("🏠 Home", key="home_btn", use_container_width=True):
        st.session_state.current_note = None
        st.session_state.show_mode_selector = False

    st.divider()
    
    # Gemini API Configuration
    with st.expander("🔑 Gemini API", expanded=False):
        st.session_state.gemini_api_key = st.text_input(
            "Gemini API Key", 
            type="password", 
            help="Enter your Google Gemini API key",
            value=st.session_state.get('gemini_api_key', '')
        )
        if st.session_state.gemini_api_key:
            st.success("✅ API Key configured!")
        else:
            st.warning("⚠️ Add API key for full functionality")
    
    st.subheader("📝 Your Notes")
    if st.button("➕ New Analysis", use_container_width=True):
        st.session_state.show_mode_selector = True
        st.session_state.current_note = None

    # Display notes
    for note in st.session_state.notes:
        doc_count = len(note.get('documents', []))
        mode_icon = {"Research": "🔬", "Study": "📚", "Normal": "📝"}.get(note['mode'], "📝")
        
        if st.button(
            f"{mode_icon} {note['title']}\n{note['mode']} • {doc_count} docs • {note['date']}", 
            key=f"note_{note['id']}", 
            use_container_width=True
        ):
            st.session_state.current_note = note
            st.session_state.show_mode_selector = False

    # Statistics
    if st.session_state.notes:
        st.divider()
        st.subheader("📊 Statistics")
        total_docs = sum(len(note.get('documents', [])) for note in st.session_state.notes)
        total_notes = len(st.session_state.notes)
        st.metric("Total Notes", total_notes)
        st.metric("Total Documents", total_docs)

# =====================
# MAIN CONTENT
# =====================
st.markdown(custom_css, unsafe_allow_html=True)

if st.session_state.show_mode_selector:
    st.title("🤖 Create New AI Analysis")
    st.markdown("Choose a mode and start analyzing your documents with AI!")
    
    note_title = st.text_input("Analysis Title", placeholder="e.g., 'Research Paper Analysis' or 'Study Notes Review'")
    
    st.subheader("Select Analysis Mode")
    mode_cols = st.columns(3)
    modes = {
        "Research": "🔬 Advanced research analysis with side-by-side layout",
        "Study": "📚 Study notes with Q&A and summaries", 
        "Normal": "📝 General document analysis and chat"
    }
    
    for i, (mode, desc) in enumerate(modes.items()):
        with mode_cols[i]:
            if st.button(f"**{mode}**\n{desc}", use_container_width=True, key=f"mode_{mode}"):
                if note_title.strip():
                    create_new_note(note_title.strip(), mode)
                    st.rerun()
                else:
                    st.error("Please enter a title for your analysis")

elif st.session_state.current_note:
    note = st.session_state.current_note
    if note['mode'] == 'Research':
        show_research_mode(note)
    else:
        show_regular_note_content(note)

else:
    # Welcome screen
    st.title("🤖 Welcome to Intituas AI")
    st.markdown("""
    ### Powerful AI-Powered Document Analysis
    *Powered by Google Gemini*
    
    **Supported File Types:**
    """)
    
    cols = st.columns(4)
    for i, (ext, desc) in enumerate(SUPPORTED_FILE_TYPES.items()):
        with cols[i % 4]:
            st.markdown(f"**{desc}**\n`.{ext}`")
    
    st.markdown("""
    ---
    ### 🚀 Features:
    - **📄 Multi-format Support**: PDF, DOCX, PPTX, Images, Text, Markdown
    - **🤖 AI Chat**: Ask questions about your documents using Gemini
    - **📝 Smart Summaries**: Generate comprehensive summaries
    - **🔬 Research Mode**: NotebookLM-style side-by-side interface
    - **🧠 Mind Maps**: Visualize document concepts
    - **🧩 Clustering**: Group similar documents
    - **💡 Smart Prompts**: Get suggested questions
    """)
    
    # Display existing notes
    if st.session_state.notes:
        st.subheader("📚 Recent Analyses")
        cols = st.columns(3)
        for i, note in enumerate(st.session_state.notes[-6:]):  # Show last 6 notes
            with cols[i % 3]:
                doc_count = len(note.get('documents', []))
                mode_icon = {"Research": "🔬", "Study": "📚", "Normal": "📝"}.get(note['mode'], "📝")
                
                if st.button(
                    f"{mode_icon} **{note['title']}**\n{note['mode']} • {doc_count} docs\n{note['date']}", 
                    key=f"card_{note['id']}", 
                    use_container_width=True
                ):
                    st.session_state.current_note = note
                    st.session_state.show_mode_selector = False
                    st.rerun()
    else:
        st.info("👆 Click '➕ New Analysis' to get started!")

# Add JavaScript for Enter key support
components.html("""
<script>
document.addEventListener('keydown', function(event) {
    if (event.key === 'Enter' && event.target.tagName === 'INPUT' && event.target.type === 'text') {
        // Find the nearest Send button and click it
        const sendButtons = document.querySelectorAll('button');
        for (let button of sendButtons) {
            if (button.textContent.includes('Send')) {
                button.click();
                break;
            }
        }
    }
});
</script>
""", height=0)

# =====================
# FOOTER
# =====================
st.markdown("---")
st.markdown(
    """
    <div style='text-align: center; color: #666;'>
        🤖 Intituas AI • Hackathon Project • Powered by Google Gemini
    </div>
    """, 
    unsafe_allow_html=True
)